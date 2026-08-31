# -*- coding: utf-8 -*-
"""
한의원 다이어트 전-후 비교 PPT 자동 생성기
- 표준모드: 구도당 1슬라이드(전-후 1세트) = 16슬라이드 + 체성분표 1슬라이드 = 17슬라이드
- 장기모드: 구도당 2슬라이드(세트1, 세트2를 이어서 배치) = 32슬라이드 + 체성분표 1슬라이드 = 33슬라이드
- 슬라이드 순서: 구도1(세트1) -> 구도1(세트2) -> 구도2(세트1) -> ... (사용자 확정 방식)
"""
import os
from datetime import date
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- 16개 구도 정의 (번호, 라벨) - 최종 참고 이미지 기준 ----
COMPOS = [
    (1, "전면_전신(1)"),
    (2, "전면_전신(2)"),
    (3, "전면_체간"),
    (4, "전면_상반신(1)"),
    (5, "전면_상반신(2)"),
    (6, "전면_하반신"),
    (7, "측면_전신"),
    (8, "측면_체간"),
    (9, "측면_상반신"),
    (10, "측면_하반신"),
    (11, "후면_전신(1)"),
    (12, "후면_전신(2)"),
    (13, "후면_체간"),
    (14, "후면_상반신(1)"),
    (15, "후면_상반신(2)"),
    (16, "후면_하반신"),
]
WIDE_COMPOS = {5, 15}  # 유일하게 가로로 긴 구도

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DARK = RGBColor(0x28, 0x28, 0x26)
MUTED = RGBColor(0x70, 0x70, 0x6a)
LINE = RGBColor(0xC8, 0xC6, 0xC0)


def _fit_box(img_path, box_w, box_h):
    """이미지 원본 비율을 유지하며 box 안에 맞는 (w,h) 계산"""
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    return int(iw * scale), int(ih * scale)


def _add_guide_tag(slide, text):
    """좌측상단 구도 텍스트 박스"""
    box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(6.5), Inches(0.5))
    tf = box.text_frame
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = DARK
    # 테두리 박스
    box.line.color.rgb = LINE
    box.line.width = Pt(1)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _add_photo(slide, path, cx, cy, box_w, box_h, tag=None, date_text=None):
    """사진을 지정 영역 중앙에 비율유지 배치. tag=전/후 라벨, date_text=날짜 캡션"""
    w, h = _fit_box(path, box_w, box_h)
    x = cx + (box_w - w) / 2
    y = cy + (box_h - h) / 2
    slide.shapes.add_picture(path, x, y, width=w, height=h)

    if tag:
        tb = slide.shapes.add_textbox(x + Inches(0.05), y + Inches(0.05), Inches(0.6), Inches(0.35))
        tf = tb.text_frame
        tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        run = tf.paragraphs[0].add_run()
        run.text = tag
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = DARK
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tb.line.color.rgb = LINE
        tb.line.width = Pt(0.75)

    if date_text:
        dtb = slide.shapes.add_textbox(x, y + h + Inches(0.05), w, Inches(0.35))
        tf = dtb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = date_text
        run.font.size = Pt(13)
        run.font.color.rgb = DARK


def _add_photo_at(slide, path, x, y, w, h, tag=None, date_text=None):
    """이미 계산된 (x,y,w,h) 그대로 사진 배치 (비율 유지는 호출부에서 처리됨)"""
    slide.shapes.add_picture(path, x, y, width=w, height=h)

    if tag:
        tb = slide.shapes.add_textbox(x + Inches(0.05), y + Inches(0.05), Inches(0.6), Inches(0.35))
        tf = tb.text_frame
        tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        run = tf.paragraphs[0].add_run()
        run.text = tag
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = DARK
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tb.line.color.rgb = LINE
        tb.line.width = Pt(0.75)

    if date_text:
        dtb = slide.shapes.add_textbox(x, y + h + Inches(0.05), w, Inches(0.35))
        tf = dtb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = date_text
        run.font.size = Pt(13)
        run.font.color.rgb = DARK


def add_photo_slide(prs, num, label, before_path, after_path, before_date=None, after_date=None,
                     show_dates=False, wide=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    _add_guide_tag(slide, f"{num}. {label}")

    top_margin = Inches(1.0)
    bottom_margin = Inches(0.9) if show_dates else Inches(0.5)
    side_margin = Inches(0.5)
    gap = Inches(0.2)  # 사진과 사진 사이 실제 간격

    avail_w = SLIDE_W - side_margin * 2
    avail_h = SLIDE_H - top_margin - bottom_margin

    if wide:
        # 5, 15번(팔벌림): 대각선 배치 - 전 사진은 좌측상단, 후 사진은 우측하단
        # 대각선 배치가 뚜렷이 보이도록 세로 공간의 72%만 사용 (겹치지 않고 어긋나게)
        diag_box_w = avail_w * 0.52
        diag_box_h = avail_h * 0.72

        bw, bh = _fit_box(before_path, diag_box_w, diag_box_h)
        aw, ah = _fit_box(after_path, diag_box_w, diag_box_h)

        before_x = side_margin
        before_y = top_margin
        after_x = side_margin + avail_w - aw
        after_y = top_margin + avail_h - ah

        _add_photo_at(slide, before_path, before_x, before_y, bw, bh,
                      tag="전", date_text=before_date if show_dates else None)
        _add_photo_at(slide, after_path, after_x, after_y, aw, ah,
                      tag="후", date_text=after_date if show_dates else None)
        return slide

    # 1) 각 사진을 세로 공간(avail_h) 기준으로 먼저 비율에 맞게 실제 크기 계산
    #    (박스를 크게 잡고 안에서 비율유지 배치하면 사진 비율에 따라 빈 여백이 생기므로,
    #     실제 렌더 크기를 먼저 구해서 그 크기끼리 딱 붙여 배치한다)
    bw, bh = _fit_box(before_path, avail_w, avail_h)
    aw, ah = _fit_box(after_path, avail_w, avail_h)

    total_w = bw + gap + aw
    if total_w > avail_w:
        # 혹시 둘을 합친 폭이 넘치면 동일 비율로 축소
        scale = avail_w / total_w
        bw, bh = int(bw * scale), int(bh * scale)
        aw, ah = int(aw * scale), int(ah * scale)
        total_w = bw + gap + aw

    start_x = side_margin + (avail_w - total_w) / 2
    before_x = start_x
    after_x = start_x + bw + gap

    before_y = top_margin + (avail_h - bh) / 2
    after_y = top_margin + (avail_h - ah) / 2

    _add_photo_at(slide, before_path, before_x, before_y, bw, bh,
                  tag="전", date_text=before_date if show_dates else None)
    _add_photo_at(slide, after_path, after_x, after_y, aw, ah,
                  tag="후", date_text=after_date if show_dates else None)

    return slide


import re

RED = RGBColor(0xC8, 0x30, 0x30)
TARGET_FILL = RGBColor(0xF7, 0xDD, 0xEC)  # 목표치 컬럼 연분홍 배경
HEADER_FILL = RGBColor(0xD9, 0xEE, 0xF5)  # 헤더 연하늘 배경


def _parse_value(text):
    """'67.3kg' -> (67.3, 'kg'). 숫자를 못찾으면 (None, 원문)."""
    if text is None or text == "":
        return None, ""
    m = re.match(r"^\s*(-?\d+(?:\.\d+)?)\s*(.*)$", str(text))
    if not m:
        return None, str(text)
    return float(m.group(1)), m.group(2).strip()


def _format_change(start_text, end_text):
    """시작/마지막 값으로 변화량 문자열 계산. 예: '67.3kg','59.2kg' -> '-8.1kg'"""
    sv, unit = _parse_value(start_text)
    ev, _ = _parse_value(end_text)
    if sv is None or ev is None:
        return "-"
    diff = round(ev - sv, 2)
    sign = "+" if diff > 0 else ("" if diff < 0 else "±")
    # 정수면 소수점 생략
    diff_str = f"{diff:g}"
    return f"{sign}{diff_str}{unit}"


def add_body_comp_slide(prs, patient_name, rows, mode):
    """
    체성분 검사 변화 표 슬라이드 (마지막 슬라이드)
    rows: [{"label":str, "start":str, "mid":str|None, "end":str, "target":str, "highlight":bool}, ...]
    - highlight=True 인 행은 항목명+수치를 빨간 글씨로 표시
    - 변화량은 start/end 값으로 자동 계산
    - 목표치 컬럼은 연분홍 배경
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "체성분 검사 변화"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = DARK

    cols = ["항목", "시작", "중간", "마지막", "변화량", "목표치(적정치)"]
    n_rows = len(rows) + 1
    n_cols = len(cols)

    table_w = SLIDE_W - Inches(1.0)
    table_h = Inches(0.42) * n_rows
    table_x = Inches(0.5)
    table_y = Inches(1.0)

    gframe = slide.shapes.add_table(n_rows, n_cols, table_x, table_y, table_w, table_h)
    table = gframe.table

    for c, colname in enumerate(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        tf_c = cell.text_frame
        tf_c.paragraphs[0].text = colname
        run = tf_c.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = DARK

    for r, row in enumerate(rows, start=1):
        label = row["label"]
        start_v = row.get("start", "")
        mid_v = row.get("mid") or ""
        end_v = row.get("end", "")
        target_v = row.get("target", "")
        highlight = row.get("highlight", False)
        change_v = _format_change(start_v, end_v)

        values = [label, start_v, mid_v, end_v, change_v, target_v]
        for c, val in enumerate(values):
            cell = table.cell(r, c)
            tf_c = cell.text_frame
            tf_c.paragraphs[0].text = str(val) if val else ""
            if tf_c.paragraphs[0].runs:
                run = tf_c.paragraphs[0].runs[0]
                run.font.size = Pt(12.5)
                run.font.color.rgb = RED if highlight else DARK
                if c == 4:  # 변화량 컬럼은 항상 굵게 강조
                    run.font.bold = True
            if c == 5:  # 목표치 컬럼 배경
                cell.fill.solid()
                cell.fill.fore_color.rgb = TARGET_FILL

    return slide


def build_presentation(patient_name, mode, photo_sets, body_comp_rows, output_path):
    """
    photo_sets: { comp_num: [ {before, after, before_date, after_date}, ... ] }
                표준모드는 리스트 길이 1, 장기모드는 길이 2 (세트1, 세트2)
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    first_slide_done = False
    for num, label in COMPOS:
        sets = photo_sets.get(num, [])
        wide = num in WIDE_COMPOS
        for s in sets:
            show_dates = not first_slide_done  # 전체 덱의 첫 슬라이드에만 날짜 표기
            add_photo_slide(
                prs, num, label,
                s["before"], s["after"],
                before_date=s.get("before_date"), after_date=s.get("after_date"),
                show_dates=show_dates, wide=wide,
            )
            first_slide_done = True

    add_body_comp_slide(prs, patient_name, body_comp_rows, mode)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    CROP_DIR = "/home/claude/cropped"
    LABEL_TO_FILE = {
        1: "전면_전신_다리벌림", 2: "전면_전신_다리오므림", 3: "전면_체간몸통", 4: "전면_상반신",
        5: "전면_상반신_팔벌림", 6: "전면_하반신_앞",
        7: "측면_전신_우측", 8: "측면_체간_우측", 9: "측면_우측_상반신", 10: "측면_우측_하반신",
        11: "뒷면_전신_다리벌림", 12: "뒷면_전신_다리오므림", 13: "뒷면_체간몸통", 14: "뒷면_상반신",
        15: "뒷면_상반신_팔벌림", 16: "뒷면_하반신_뒤",
    }

    photo_sets = {}
    for num in range(1, 17):
        fname = LABEL_TO_FILE[num]
        before = os.path.join(CROP_DIR, f"before_{num:02d}_{fname}.png")
        after = os.path.join(CROP_DIR, f"after_{num:02d}_{fname}.png")
        photo_sets[num] = [{
            "before": before, "after": after,
            "before_date": "2026.01.05", "after_date": "2026.03.10",
        }]

    body_comp_rows = [
        {"label": "체중(kg)", "start": "67.3kg", "mid": None, "end": "59.2kg", "target": "60kg(55kg)", "highlight": False},
        {"label": "신장(cm)", "start": "162.5cm", "mid": None, "end": "162.5cm", "target": "-", "highlight": False},
        {"label": "체지방량(kg)", "start": "26.0kg", "mid": None, "end": "19.7kg", "target": "12.8kg", "highlight": False},
        {"label": "골격근량(kg)", "start": "22.5kg", "mid": None, "end": "21.4kg", "target": "24.5kg", "highlight": False},
        {"label": "체지방률(%)", "start": "38.6%", "mid": None, "end": "33.2%", "target": "18~28%", "highlight": False},
        {"label": "BMI(kg/m2)", "start": "25.5", "mid": None, "end": "22.4", "target": "18.5~25", "highlight": False},
        {"label": "복부지방률", "start": "1.00", "mid": None, "end": "0.93", "target": "0.75~0.85", "highlight": False},
        {"label": "복부둘레", "start": "96.1cm", "mid": None, "end": "85.2cm", "target": "73cm", "highlight": True},
        {"label": "엉덩이둘레", "start": "96.1cm", "mid": None, "end": "91.7cm", "target": "89.2cm", "highlight": True},
        {"label": "가슴둘레", "start": "95.0cm", "mid": None, "end": "89.2cm", "target": "84.1cm", "highlight": True},
        {"label": "우측상완둘레", "start": "32.3cm", "mid": None, "end": "29.6cm", "target": "26.6cm", "highlight": False},
        {"label": "우측대퇴둘레", "start": "50.6cm", "mid": None, "end": "48.4cm", "target": "48.6cm", "highlight": False},
    ]

    out = build_presentation(
        patient_name="홍길동",
        mode="standard",
        photo_sets=photo_sets,
        body_comp_rows=body_comp_rows,
        output_path="/home/claude/ppt_gen/홍길동님_프레젠테이션_20260830.pptx",
    )
    print("생성완료:", out)
