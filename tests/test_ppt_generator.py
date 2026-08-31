# -*- coding: utf-8 -*-
"""PPT 생성 로직 검증. 실사진 없이 PIL로 만든 더미 이미지로 슬라이드 수/배치를 확인한다."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ppt_generator.compos import COMPOS, WIDE_COMPOS
from ppt_generator.generate_ppt import _format_change, build_presentation

BODY_COMP_ROWS = [
    {"label": "체중", "start": "70kg", "mid": None, "end": "60kg", "target": "55kg", "highlight": True},
]


def _photo_sets_for(n_sets: int, make_dummy_photo) -> dict[int, list[dict]]:
    photo_sets: dict[int, list[dict]] = {}
    for num, _label in COMPOS:
        size = (960, 540) if num in WIDE_COMPOS else (720, 960)
        sets = []
        for s in range(1, n_sets + 1):
            before = make_dummy_photo(f"before_{num}_{s}.png", size=size)
            after = make_dummy_photo(f"after_{num}_{s}.png", size=size)
            sets.append({"before": before, "after": after,
                         "before_date": "2026.01.01", "after_date": "2026.02.01"})
        photo_sets[num] = sets
    return photo_sets


def test_standard_mode_slide_count(tmp_path, make_dummy_photo):
    photo_sets = _photo_sets_for(1, make_dummy_photo)
    out = build_presentation("홍길동", "standard", photo_sets, BODY_COMP_ROWS, str(tmp_path / "out.pptx"))

    prs = Presentation(out)
    assert len(prs.slides) == 17  # 16개 구도 + 체성분표 1


def test_long_mode_slide_count(tmp_path, make_dummy_photo):
    photo_sets = _photo_sets_for(2, make_dummy_photo)
    out = build_presentation("홍길동", "long", photo_sets, BODY_COMP_ROWS, str(tmp_path / "out_long.pptx"))

    prs = Presentation(out)
    assert len(prs.slides) == 33  # 32개(구도x2세트) + 체성분표 1


def test_partial_deck_skips_missing_compos(tmp_path, make_dummy_photo):
    """일부 구도 사진이 없어도 전체 생성이 실패하지 않고 해당 구도만 빠진다."""
    photo_sets = _photo_sets_for(1, make_dummy_photo)
    del photo_sets[3]  # 구도 3번 세트를 통째로 제거

    out = build_presentation("홍길동", "standard", photo_sets, BODY_COMP_ROWS, str(tmp_path / "out_partial.pptx"))
    prs = Presentation(out)
    assert len(prs.slides) == 16  # 15개 구도 + 체성분표 1


def test_wide_compos_diagonal_non_overlapping(tmp_path, make_dummy_photo):
    photo_sets = _photo_sets_for(1, make_dummy_photo)
    out = build_presentation("홍길동", "standard", photo_sets, BODY_COMP_ROWS, str(tmp_path / "out_wide.pptx"))

    prs = Presentation(out)
    wide_slide = prs.slides[4]  # COMPOS 순서상 5번 구도가 5번째(0-indexed 4) 슬라이드
    pictures = [sh for sh in wide_slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 2

    p1, p2 = pictures
    # 대각선 배치: 전 사진(좌상단) / 후 사진(우하단) - 좌표가 같은 방향으로 어긋나야 한다
    diagonal = (p1.left < p2.left and p1.top < p2.top) or (p2.left < p1.left and p2.top < p1.top)
    assert diagonal


def test_only_first_slide_has_date_caption(tmp_path, make_dummy_photo):
    photo_sets = _photo_sets_for(1, make_dummy_photo)
    out = build_presentation("홍길동", "standard", photo_sets, BODY_COMP_ROWS, str(tmp_path / "out_date.pptx"))

    prs = Presentation(out)

    def textbox_texts(slide):
        return [
            run.text
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        ]

    first_texts = textbox_texts(prs.slides[0])
    second_texts = textbox_texts(prs.slides[1])
    assert any("2026.01.01" in t for t in first_texts)
    assert not any("2026.01.01" in t for t in second_texts)


def test_body_comp_table_highlight_and_change_column(tmp_path, make_dummy_photo):
    photo_sets = _photo_sets_for(1, make_dummy_photo)
    rows = [
        {"label": "복부둘레", "start": "96.1cm", "mid": None, "end": "85.2cm", "target": "73cm", "highlight": True},
        {"label": "신장", "start": "162.5cm", "mid": None, "end": "162.5cm", "target": "-", "highlight": False},
    ]
    out = build_presentation("홍길동", "standard", photo_sets, rows, str(tmp_path / "out_table.pptx"))

    prs = Presentation(out)
    table_slide = prs.slides[-1]
    table = next(sh.table for sh in table_slide.shapes if sh.has_table)

    # 헤더(0행) + 2개 데이터 행
    assert len(table.rows) == 3
    # 변화량 컬럼(4번 인덱스)이 자동 계산되어 표시됨
    assert table.cell(1, 4).text_frame.text == "-10.9cm"
    assert table.cell(2, 4).text_frame.text == "±0cm"


def test_format_change_examples():
    assert _format_change("67.3kg", "59.2kg") == "-8.1kg"
    assert _format_change("162.5cm", "162.5cm") == "±0cm"
    assert _format_change("", "60kg") == "-"
    assert _format_change("20", "25") == "+5"
